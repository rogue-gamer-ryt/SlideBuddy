from pptx import Presentation


def create_presentation(outline, introduction, main_topics, conclusion):
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = outline['title']
    subtitle.text = "Subtitle here"

    # Introduction slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title.text = "Introduction"
    content_placeholder.text = introduction

    # Main topics slides
    for topic in main_topics:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        title.text = topic['title']
        content_placeholder.text = topic['content']

    # Conclusion slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title.text = "Conclusion"
    content_placeholder.text = conclusion

    prs.save('presentation.pptx')
    return 'presentation.pptx'


# Example usage with generated content
outline = {
    'title': "Your Presentation Title",
    'main_topics': [
        {'title': "Main Topic 1"},
        {'title': "Main Topic 2"},
        {'title': "Main Topic 3"},
    ]
}

introduction = "This is the introduction."
main_topics = [
    {'title': "Main Topic 1", 'content': "Content for main topic 1."},
    {'title': "Main Topic 2", 'content': "Content for main topic 2."},
    {'title': "Main Topic 3", 'content': "Content for main topic 3."},
]
conclusion = "This is the conclusion."

create_presentation(outline, introduction, main_topics, conclusion)
